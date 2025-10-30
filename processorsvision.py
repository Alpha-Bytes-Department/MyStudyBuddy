import os
import base64
from pathlib import Path
from typing import Dict, List, Any
import openai
from openai import OpenAI
import json
import streamlit as st
import tempfile
from dotenv import load_dotenv
import re

# Load environment variables
load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# Document processing libraries
from docx import Document
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.table import Table
from docx.text.paragraph import Paragraph
import PyPDF2
from pdf2image import convert_from_path
from pptx import Presentation
from PIL import Image
import io

class DocumentExtractor:
    def __init__(self):
        """Initialize Document Extractor"""
        pass
    
    def remove_numbering(self, text: str) -> str:
        """Remove numbering patterns from text"""
        if not text:
            return text
        
        # Split into lines
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            # Remove patterns like "1. ", "2. ", "10. ", "1: ", etc. at start of line
            cleaned = re.sub(r'^\s*\d+[\.\:\)]\s*', '', line)
            # Remove patterns like "0 : ", "1 : " with spaces
            cleaned = re.sub(r'^\s*\d+\s*:\s*', '', cleaned)
            cleaned_lines.append(cleaned)
        
        return '\n'.join(cleaned_lines)
        
    def encode_image_to_base64(self, image_path: str) -> str:
        """Convert image to base64 string"""
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')
    
    def image_to_base64(self, image: Image.Image) -> str:
        """Convert PIL Image to base64"""
        buffered = io.BytesIO()
        image.save(buffered, format="PNG")
        return base64.b64encode(buffered.getvalue()).decode('utf-8')
    
    def extract_with_vision(self, image_path: str = None, image_base64: str = None) -> str:
        """Extract content from image using GPT-4 Vision"""
        try:
            if image_path:
                image_base64 = self.encode_image_to_base64(image_path)
            
            response = client.chat.completions.create(
                model="gpt-4o",  # or "gpt-4-vision-preview"
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": """Extract ALL content from this image including:
- All text (printed and handwritten)
- All diagrams, charts, graphs - describe them in detail
- All tables with their exact structure
- Mathematical equations and formulas
- Any annotations, notes, or markings


Preserve the original structure and formatting as much as possible.
If there are diagrams/images, provide detailed descriptions including:
- Type of diagram (flowchart, graph, chart, etc.)
- All labels, legends, and annotations
- Relationships between elements
- Colors, shapes, and visual elements
- Data points and values

DO NOT use numbered lists or bullet point numbering. Present information naturally without line numbers."""
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{image_base64}",
                                    "detail": "high"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=4096
            )
            
            extracted_text = response.choices[0].message.content
            # Apply numbering removal as additional safeguard
            return self.remove_numbering(extracted_text)
            
        except Exception as e:
            return f"Error extracting image: {str(e)}"
    
    def extract_from_docx(self, file_path: str) -> Dict[str, Any]:
        """Extract everything from DOCX file"""
        doc = Document(file_path)
        content = {
            "text": "",
            "tables": [],
            "images": []
        }
        
        # Extract text and images in order
        text_items = []
        for element in doc.element.body:
            if isinstance(element, CT_P):
                para = Paragraph(element, doc)
                if para.text.strip():
                    text_items.append(para.text)
                
                # Check for images in paragraph
                for run in para.runs:
                    if run._element.xpath('.//a:blip'):
                        for rel in run.part.rels.values():
                            if "image" in rel.target_ref:
                                try:
                                    image_data = rel.target_part.blob
                                    image = Image.open(io.BytesIO(image_data))
                                    image_base64 = self.image_to_base64(image)
                                    extracted = self.extract_with_vision(image_base64=image_base64)
                                    content["images"].append({
                                        "type": "embedded_image",
                                        "extracted_content": extracted
                                    })
                                except Exception as e:
                                    content["images"].append({
                                        "type": "embedded_image",
                                        "error": str(e)
                                    })
            
            elif isinstance(element, CT_Tbl):
                table = Table(element, doc)
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                content["tables"].append(table_data)
        
        # Join all text items into single string
        content["text"] = "\n\n".join(text_items)
        return content
    
    def extract_from_pdf(self, file_path: str) -> Dict[str, Any]:
        """Extract everything from PDF file"""
        content = {
            "text": "",
        }
        
        # Extract text
        text_pages = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    text_pages.append(f"--- Page {page_num + 1} ---\n{text}")
        except Exception as e:
            text_pages.append(f"Text extraction error: {str(e)}")
        
        # Join all pages into single string
        content["text"] = "\n\n".join(text_pages)
        return content
    
    def extract_from_pptx(self, file_path: str) -> Dict[str, Any]:
        """Extract everything from PPTX file"""
        prs = Presentation(file_path)
        content = {
            "slides": []
        }
        
        for slide_num, slide in enumerate(prs.slides):
            slide_content = {
                "slide_number": slide_num + 1,
                "text": "",
                "tables": [],
                "images": []
            }
            
            # Extract text
            text_items = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_items.append(shape.text)
                
                # Extract tables
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    slide_content["tables"].append(table_data)
                
                # Extract images
                if shape.shape_type == 13:  # Picture
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        pil_image = Image.open(io.BytesIO(image_bytes))
                        image_base64 = self.image_to_base64(pil_image)
                        extracted = self.extract_with_vision(image_base64=image_base64)
                        slide_content["images"].append({
                            "type": "embedded_image",
                            "extracted_content": extracted
                        })
                    except Exception as e:
                        slide_content["images"].append({
                            "type": "embedded_image",
                            "error": str(e)
                        })
            
            # Join all text items into single string
            slide_content["text"] = "\n".join(text_items)
            content["slides"].append(slide_content)
        
        return content
    
    def extract_from_image(self, file_path: str) -> Dict[str, Any]:
        """Extract everything from image file (including handwritten notes)"""
        st.write(f"Extracting from image: {Path(file_path).name}")
        extracted_content = self.extract_with_vision(image_path=file_path)
        
        return {
            "file_type": "image",
            "extracted_content": extracted_content
        }
    
    def extract_from_file(self, file_path: str) -> Dict[str, Any]:
        """Main method to extract from any supported file type"""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        st.write(f"Processing: {file_path.name}")
        
        if extension == '.docx':
            return {
                "file_name": file_path.name,
                "file_type": "docx",
                "content": self.extract_from_docx(str(file_path))
            }
        
        elif extension == '.pdf':
            return {
                "file_name": file_path.name,
                "file_type": "pdf",
                "content": self.extract_from_pdf(str(file_path))
            }
        
        elif extension in ['.ppt', '.pptx']:
            return {
                "file_name": file_path.name,
                "file_type": "pptx",
                "content": self.extract_from_pptx(str(file_path))
            }
        
        elif extension in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.webp']:
            return {
                "file_name": file_path.name,
                "file_type": "image",
                "content": self.extract_from_image(str(file_path))
            }
        
        else:
            return {
                "file_name": file_path.name,
                "error": f"Unsupported file type: {extension}"
            }
    
    def save_extraction(self, extracted_data: Dict[str, Any], output_path: str):
        """Save extracted content to JSON file"""
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(extracted_data, f, indent=2, ensure_ascii=False)
        st.success(f"Extraction saved to: {output_path}")


# Streamlit Display Functions
def display_docx_content(content):
    """Display DOCX extracted content"""
    st.subheader("📝 Text Content")
    if content.get('text'):
        st.text(content['text'])
    
    if content.get('tables'):
        st.subheader("📊 Tables")
        for idx, table in enumerate(content['tables'], 1):
            st.write(f"**Table {idx}:**")
            st.table(table)
    
    if content.get('images'):
        st.subheader("🖼️ Images")
        for idx, img in enumerate(content['images'], 1):
            with st.expander(f"Image {idx}", expanded=True):
                if 'extracted_content' in img:
                    st.markdown(img['extracted_content'])
                else:
                    st.error(img.get('error', 'Unknown error'))

def display_pdf_content(content):
    """Display PDF extracted content"""
    st.subheader("📄 Text Extraction")
    if content.get('text'):
        st.text(content['text'])

def display_pptx_content(content):
    """Display PPTX extracted content"""
    for slide in content.get('slides', []):
        slide_num = slide.get('slide_number', 'Unknown')
        st.subheader(f"🎯 Slide {slide_num}")
        
        col1, col2 = st.columns([2, 1])
        
        with col1:
            if slide.get('text'):
                st.markdown("**Text:**")
                st.text(slide['text'])
            
            if slide.get('tables'):
                st.markdown("**Tables:**")
                for idx, table in enumerate(slide['tables'], 1):
                    st.write(f"Table {idx}:")
                    st.table(table)
        
        with col2:
            if slide.get('images'):
                st.markdown("**Images:**")
                for idx, img in enumerate(slide['images'], 1):
                    with st.expander(f"Image {idx}", expanded=False):
                        if 'extracted_content' in img:
                            st.markdown(img['extracted_content'])
                        else:
                            st.error(img.get('error', 'Unknown error'))
        
        st.divider()

def display_image_content(content):
    """Display image extracted content"""
    st.subheader("🖼️ Extracted Content")
    if 'extracted_content' in content:
        st.markdown(content['extracted_content'])
    else:
        st.error("No content extracted")


# Main Streamlit App
def main():
    st.set_page_config(
        page_title="Document Content Extractor",
        page_icon="📄",
        layout="wide"
    )
    
    st.title("📄 Document Content Extractor")
    st.markdown("Extract text, tables, diagrams, and handwritten notes from documents using OpenAI GPT-4 Vision")
    
    # Initialize session state
    if 'extracted_data' not in st.session_state:
        st.session_state.extracted_data = None
    
    # File uploader
    uploaded_file = st.file_uploader(
        "Upload your document",
        type=['docx', 'pdf', 'pptx', 'ppt', 'png', 'jpg', 'jpeg', 'bmp', 'tiff', 'webp'],
        help="Upload DOCX, PDF, PPTX, or Image files"
    )
    
    if uploaded_file:
        # Display original filename
        st.info(f"📎 Uploaded: **{uploaded_file.name}**")
        
        # Save uploaded file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix=Path(uploaded_file.name).suffix) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_path = tmp_file.name
        
        # Extract button
        if st.button("🚀 Extract Content", type="primary", use_container_width=True):
            with st.spinner("Extracting content... This may take a few moments..."):
                try:
                    extractor = DocumentExtractor()
                    result = extractor.extract_from_file(tmp_path)
                    # Override with original filename
                    result['file_name'] = uploaded_file.name
                    st.session_state.extracted_data = result
                    st.success("✅ Extraction completed successfully!")
                except Exception as e:
                    st.error(f"❌ Error: {str(e)}")
                    st.session_state.extracted_data = None
    
    # Display extracted content
    if st.session_state.extracted_data:
        result = st.session_state.extracted_data
        
        st.divider()
        
        # File info
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("File Name", result.get('file_name', 'Unknown'))
        with col2:
            st.metric("File Type", result.get('file_type', 'Unknown').upper())
        with col3:
            json_str = json.dumps(result, indent=2, ensure_ascii=False)
            st.download_button(
                label="💾 Download JSON",
                data=json_str,
                file_name=f"{Path(result.get('file_name', 'extraction')).stem}_extracted.json",
                mime="application/json"
            )
        
        st.divider()
        
        # Display content based on file type
        content = result.get('content', {})
        file_type = result.get('file_type', '').lower()
        
        if file_type == 'docx':
            display_docx_content(content)
        elif file_type == 'pdf':
            display_pdf_content(content)
        elif file_type in ['pptx', 'ppt']:
            display_pptx_content(content)
        elif file_type == 'image':
            display_image_content(content)
        else:
            st.error(result.get('error', 'Unknown error'))
        
        # Raw JSON view
        with st.expander("🔍 View Raw JSON Data"):
            st.json(result)


if __name__ == "__main__":
    main()