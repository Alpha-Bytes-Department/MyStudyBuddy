import streamlit as st
import pytesseract
from PIL import Image
import pdf2image
import docx
import io
import os
from pptx import Presentation
import tempfile
import cv2
import numpy as np

# Configure Tesseract path (uncomment and modify if needed)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'  # Windows
# pytesseract.pytesseract.tesseract_cmd = r'/usr/bin/tesseract'  # Linux/Mac

class DocumentProcessor:
    """Process various document types and extract text using Tesseract OCR"""
    
    @staticmethod
    def preprocess_image(image):
        """Preprocess image for better OCR accuracy"""
        # Convert PIL Image to numpy array
        img_array = np.array(image)
        
        # Convert to grayscale
        gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
        
        # Apply thresholding to get binary image
        _, binary = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        
        # Denoise
        denoised = cv2.fastNlMeansDenoising(binary, None, 10, 7, 21)
        
        # Convert back to PIL Image
        return Image.fromarray(denoised)
    
    @staticmethod
    def extract_text_from_image(image, preprocess=True):
        """Extract text from image using Tesseract OCR"""
        try:
            if preprocess:
                image = DocumentProcessor.preprocess_image(image)
            
            # Use custom config for better accuracy
            custom_config = r'--oem 3 --psm 6'
            text = pytesseract.image_to_string(image, config=custom_config)
            return text.strip()
        except Exception as e:
            return f"Error extracting text from image: {str(e)}"
    
    @staticmethod
    def extract_text_from_pdf(pdf_file, use_ocr=True):
        """Extract text from PDF (with OCR support for scanned PDFs)"""
        try:
            # Save uploaded file temporarily
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
                tmp_file.write(pdf_file.read())
                tmp_path = tmp_file.name
            
            # Convert PDF to images
            images = pdf2image.convert_from_path(tmp_path, dpi=300)
            
            extracted_text = []
            
            for i, image in enumerate(images):
                st.write(f"Processing page {i+1}/{len(images)}...")
                
                if use_ocr:
                    # Use OCR for each page
                    page_text = DocumentProcessor.extract_text_from_image(image)
                else:
                    # Simple OCR without preprocessing
                    page_text = pytesseract.image_to_string(image)
                
                if page_text.strip():
                    extracted_text.append(f"--- Page {i+1} ---\n{page_text}\n")
            
            # Clean up temporary file
            os.unlink(tmp_path)
            
            return "\n".join(extracted_text)
        
        except Exception as e:
            return f"Error extracting text from PDF: {str(e)}"
    
    @staticmethod
    def extract_text_from_docx(docx_file):
        """Extract text from DOCX file"""
        try:
            doc = docx.Document(docx_file)
            text_content = []
            
            # Extract text from paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join([cell.text for cell in row.cells])
                    if row_text.strip():
                        text_content.append(row_text)
            
            return "\n".join(text_content)
        
        except Exception as e:
            return f"Error extracting text from DOCX: {str(e)}"
    
    @staticmethod
    def extract_text_from_pptx(pptx_file):
        """Extract text from PPTX file"""
        try:
            prs = Presentation(pptx_file)
            text_content = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = [f"--- Slide {i+1} ---"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    
                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = ' | '.join([cell.text for cell in row.cells])
                            if row_text.strip():
                                slide_text.append(row_text)
                
                if len(slide_text) > 1:  # If there's content beyond the header
                    text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
        
        except Exception as e:
            return f"Error extracting text from PPTX: {str(e)}"
    
    @staticmethod
    def process_file(uploaded_file, file_type, use_ocr=True):
        """Process uploaded file based on its type"""
        if file_type in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif']:
            image = Image.open(uploaded_file)
            return DocumentProcessor.extract_text_from_image(image, preprocess=use_ocr)
        
        elif file_type == 'pdf':
            return DocumentProcessor.extract_text_from_pdf(uploaded_file, use_ocr=use_ocr)
        
        elif file_type == 'docx':
            return DocumentProcessor.extract_text_from_docx(uploaded_file)
        
        elif file_type == 'pptx':
            return DocumentProcessor.extract_text_from_pptx(uploaded_file)
        
        else:
            return "Unsupported file format"


def main():
    st.set_page_config(
        page_title="Document Text Extractor",
        page_icon="📄",
        layout="wide"
    )
    
    # Custom CSS
    st.markdown("""
        <style>
        .stApp {
            max-width: 1200px;
            margin: 0 auto;
        }
        .upload-box {
            border: 2px dashed #4CAF50;
            border-radius: 10px;
            padding: 20px;
            text-align: center;
            background-color: #f9f9f9;
        }
        .extracted-text {
            background-color: #f5f5f5;
            padding: 20px;
            border-radius: 10px;
            border-left: 4px solid #4CAF50;
            max-height: 600px;
            overflow-y: auto;
            color: #000000;
            font-family: monospace;
            white-space: pre-wrap;
            word-wrap: break-word;
        }
        </style>
    """, unsafe_allow_html=True)
    
    st.title("📄 Document & Image Text Extractor")
    st.markdown("Extract text from documents (PDF, DOCX, PPTX) and images using Tesseract OCR")
    
    # Settings at the top
    use_preprocessing = True  # Always enabled for best results
    
    # Main content area
    st.markdown("### Upload Your File")
    
    uploaded_file = st.file_uploader(
        "Choose a file to extract text from",
        type=['pdf', 'docx', 'pptx', 'jpg', 'jpeg', 'png', 'bmp', 'tiff', 'gif'],
        help="Upload a document or image file"
    )
    
    if uploaded_file is not None:
        # Get file details
        file_details = {
            "Filename": uploaded_file.name,
            "File Size": f"{uploaded_file.size / 1024:.2f} KB",
            "File Type": uploaded_file.type
        }
        
        # Display file info
        col1, col2, col3 = st.columns(3)
        with col1:
            st.info(f"**📁 Filename:** {file_details['Filename']}")
        with col2:
            st.info(f"**📊 Size:** {file_details['File Size']}")
        with col3:
            st.info(f"**📄 Type:** {file_details['File Type']}")
        
        # Preview for images
        file_extension = uploaded_file.name.split('.')[-1].lower()
        if file_extension in ['jpg', 'jpeg', 'png', 'bmp', 'gif']:
            st.markdown("### 🖼️ Image Preview")
            image = Image.open(uploaded_file)
            st.image(image, caption="Uploaded Image", use_container_width=True)
            uploaded_file.seek(0)  # Reset file pointer
        
        # Extract text button
        if st.button("🔍 Extract Text", type="primary", use_container_width=True):
            with st.spinner("🔄 Processing file... This may take a moment."):
                # Process the file
                extracted_text = DocumentProcessor.process_file(
                    uploaded_file,
                    file_extension,
                    use_ocr=use_preprocessing
                )
                
                # Store in session state
                st.session_state.extracted_text = extracted_text
                st.session_state.filename = uploaded_file.name
        
        # Display extracted text
        if 'extracted_text' in st.session_state:
            st.markdown("---")
            st.markdown("### 📝 Extracted Text")
            
            # Display text in a styled container
            st.markdown(
                f'<div class="extracted-text" style="color: #000000;">{st.session_state.extracted_text.replace(chr(10), "<br>")}</div>',
                unsafe_allow_html=True
            )
            
            # Statistics
            col1, col2, col3 = st.columns(3)
            with col1:
                char_count = len(st.session_state.extracted_text)
                st.metric("📊 Characters", f"{char_count:,}")
            with col2:
                word_count = len(st.session_state.extracted_text.split())
                st.metric("📖 Words", f"{word_count:,}")
            with col3:
                line_count = len(st.session_state.extracted_text.split('\n'))
                st.metric("📄 Lines", f"{line_count:,}")
            
            # Download options
            st.markdown("---")
            
            # Download as TXT
            st.download_button(
                label="📥 Download as TXT",
                data=st.session_state.extracted_text,
                file_name=f"{st.session_state.filename}_extracted.txt",
                mime="text/plain",
                use_container_width=True
            )
    
    else:
        # Instructions when no file is uploaded
        st.markdown("""
        <div class="upload-box">
            <h3>👆 Upload a file to get started</h3>
            <p>Supports PDF, DOCX, PPTX, and various image formats</p>
        </div>
        """, unsafe_allow_html=True)
        
        # Example showcase
        st.markdown("---")
        st.markdown("### 🎯 How It Works")
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.markdown("""
            **1️⃣ Upload**
            - Select your document or image
            - Supports multiple formats
            - Drag & drop supported
            """)
        
        with col2:
            st.markdown("""
            **2️⃣ Process**
            - OCR extracts text automatically
            - Image preprocessing enhances accuracy
            - Progress shown in real-time
            """)
        
        with col3:
            st.markdown("""
            **3️⃣ Export**
            - Download as TXT file
            - Copy to clipboard
            - View statistics
            """)


if __name__ == "__main__":
    main()